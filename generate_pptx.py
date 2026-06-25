from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Couleurs ──
PRIMARY = RGBColor(0x2C, 0x3E, 0x50)
BLUE = RGBColor(0x34, 0x98, 0xDB)
BLUE_DARK = RGBColor(0x29, 0x79, 0xB3)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
MED_GRAY = RGBColor(0xBD, 0xC3, 0xC7)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
SUBTLE = RGBColor(0x7F, 0x8C, 0x8D)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)


def solid_bg(slide, color=LIGHT_GRAY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height, size=Pt(14), color=DARK_TEXT,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, items, left, top, width, height, size=Pt(16), color=DARK_TEXT, bullet_color=BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.font.name = 'Calibri'
    return txBox


def add_card(slide, title, items, left, top, width, height, header_color=BLUE):
    """Carte avec arrière-plan blanc + header coloré + contenu liste"""
    card_bg = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, WHITE)
    card_bg.shadow.inherit = False

    header = add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.05), top + Inches(0.05),
                       width - Inches(0.1), Inches(0.55), header_color)
    add_text_box(slide, title, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4),
                 Inches(0.5), size=Pt(16), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_bullet_list(slide, items, left + Inches(0.15), top + Inches(0.7),
                    width - Inches(0.3), height - Inches(0.8),
                    size=Pt(13), color=DARK_TEXT)


def add_section_header(slide, number, title, subtitle=""):
    """Bannière bleue en haut de slide"""
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.4), PRIMARY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.35), prs.slide_width, Inches(0.08), BLUE)

    add_text_box(slide, f"{number}", Inches(0.5), Inches(0.2), Inches(0.6), Inches(0.6),
                 size=Pt(28), color=BLUE, bold=True)
    add_text_box(slide, title, Inches(1.1), Inches(0.2), Inches(10), Inches(0.6),
                 size=Pt(30), color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, subtitle, Inches(1.1), Inches(0.75), Inches(10), Inches(0.4),
                     size=Pt(16), color=RGBColor(0xBD, 0xBD, 0xBD), bold=False)


def add_footer(slide, text="IT Parc TECHPARK CI  |  Silué Yacouba"):
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), prs.slide_width, Inches(0.4), PRIMARY)
    add_text_box(slide, text, Inches(0.5), Inches(7.1), Inches(12), Inches(0.4),
                 size=Pt(10), color=MED_GRAY, alignment=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, PRIMARY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15), BLUE)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), prs.slide_width, Inches(0.15), BLUE)

# Accent bar left
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.5), Inches(0.08), Inches(2.5), BLUE)

add_text_box(slide, "IT PARC TECHPARK CI", Inches(2), Inches(2.5), Inches(10), Inches(1),
             size=Pt(48), color=WHITE, bold=True)
add_text_box(slide, "Module Odoo 18 — Gestion du Parc Informatique", Inches(2), Inches(3.5), Inches(10), Inches(0.7),
             size=Pt(22), color=BLUE, bold=False)
add_text_box(slide, "Présenté par Silué Yacouba", Inches(2), Inches(4.5), Inches(10), Inches(0.5),
             size=Pt(18), color=MED_GRAY)

# Date
add_text_box(slide, "Juin 2026", Inches(2), Inches(5.2), Inches(10), Inches(0.4),
             size=Pt(14), color=MED_GRAY)

# Bottom decorative element
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), prs.slide_width, Inches(0.04), BLUE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLÉMATIQUE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, LIGHT_GRAY)
add_section_header(slide, "01", "Problématique", "Pourquoi ce projet est nécessaire")
add_footer(slide)

problems = [
    "Aucun outil centralisé → suivi manuel (Excel, papier, fichiers éparpillés)",
    "Traçabilité inexistante → on ne sait pas qui a quel équipement ni depuis quand",
    "Interventions non structurées → pas de planification, pas de suivi des coûts",
    "Contrats fournisseurs perdus → risques d'expiration sans renouvellement",
    "Aucune alerte automatique → garanties expirées, maintenances oubliées",
    "Reporting limité → pas de tableau de bord ni d'indicateurs temps réel",
]

# Left column - problem icons/cards
add_card(slide, "Problèmes identifiés", problems,
         Inches(0.8), Inches(1.8), Inches(11.5), Inches(5), RED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — OBJECTIFS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, LIGHT_GRAY)
add_section_header(slide, "02", "Objectifs du Projet", "Les 5 piliers de la solution")
add_footer(slide)

obj_data = [
    (["Plateforme unique centralisant tout le parc informatique"], GREEN),
    (["Traçabilité complète : qui, quoi, quand (affectations + historique)"], BLUE),
    (["Automatisation des alertes garanties et contrats"], ORANGE),
    (["Dashboard temps réel + exports PDF et Excel"], PURPLE),
    (["Sécurisation avec 2 profils : Technicien et Manager"], RED),
]

for i, (items, color) in enumerate(obj_data):
    y = Inches(1.8) + Inches(1.0) * i
    # Number circle
    circle = add_shape(slide, MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.1), Inches(0.5), Inches(0.5), color)
    add_text_box(slide, str(i + 1), Inches(0.6), y + Inches(0.1), Inches(0.5), Inches(0.5),
                 size=Pt(20), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Text
    add_text_box(slide, items[0], Inches(1.4), y + Inches(0.15), Inches(11), Inches(0.5),
                 size=Pt(17), color=DARK_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — PÉRIMÈTRE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, LIGHT_GRAY)
add_section_header(slide, "03", "Périmètre Fonctionnel", "9 fonctionnalités clés")
add_footer(slide)

features = [
    ("1", "Équipements", "PC, serveurs, réseau, périphériques", BLUE),
    ("2", "Affectations", "Historique employé par équipement", GREEN),
    ("3", "Interventions", "Planification, coûts, rapport PDF", ORANGE),
    ("4", "Contrats", "Fournisseurs, dates, montants FCFA", PURPLE),
    ("5", "Alertes", "Automatiques garanties + contrats", RED),
    ("6", "Import CSV", "Création masse équipements", BLUE_DARK),
    ("7", "Rapports PDF", "Fiche, intervention, contrat", GREEN),
    ("8", "Exports Excel", "Inventaire, coûts, contrats expirants", ORANGE),
    ("9", "Dashboard OWL", "4 KPI + graphique catégorie", YELLOW),
]

for i, (num, title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + col * Inches(4.1)
    y = Inches(1.7) + row * Inches(1.8)

    card_bg = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(1.5), WHITE)
    # Color top bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(3.8), Inches(0.07), color)
    # Number
    circle = add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.2), Inches(0.4), Inches(0.4), color)
    add_text_box(slide, num, x + Inches(0.15), y + Inches(0.2), Inches(0.4), Inches(0.4),
                 size=Pt(16), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, title, x + Inches(0.7), y + Inches(0.2), Inches(3), Inches(0.4),
                 size=Pt(18), color=DARK_TEXT, bold=True)
    # Desc
    add_text_box(slide, desc, x + Inches(0.7), y + Inches(0.65), Inches(3), Inches(0.6),
                 size=Pt(13), color=SUBTLE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, LIGHT_GRAY)
add_section_header(slide, "04", "Architecture Technique", "Stack technologique utilisé")
add_footer(slide)

arch_items = [
    ("Odoo 18.0", "Framework ERP complet (Python 3.11, PostgreSQL)"),
    ("Python / XML", "Modèles hérités + vues déclaratives"),
    ("OWL 2 (JS)", "Dashboard dynamique côté client"),
    ("QWeb", "Rapports PDF natifs Odoo"),
    ("xlsxwriter", "Exports Excel formatés"),
    ("SCSS", "Styles personnalisés dashboard"),
]

for i, (tech, desc) in enumerate(arch_items):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + col * Inches(4.1)
    y = Inches(1.7) + row * Inches(2.5)

    colors = [BLUE, GREEN, ORANGE, PURPLE, RED, YELLOW]
    card_bg = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.1), WHITE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(3.8), Inches(0.07), colors[i])
    # Tech name in colored box
    tech_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), y + Inches(0.3),
                         Inches(3.2), Inches(0.55), colors[i])
    add_text_box(slide, tech, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.55),
                 size=Pt(18), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, desc, x + Inches(0.3), y + Inches(1.0), Inches(3.2), Inches(0.8),
                 size=Pt(14), color=SUBTLE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — MODÈLES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, LIGHT_GRAY)
add_section_header(slide, "05", "Modèles de Données", "5 modèles interconnectés")
add_footer(slide)

model_cards = [
    ("Équipement", ["Nom, n° série, modèle", "Catégorie, site", "Valeur FCFA, garantie", "Contrat lié"], BLUE),
    ("Affectation", ["Équipement → Employé", "Département", "Dates début/fin", "Notes"], GREEN),
    ("Intervention", ["Type préventive/corrective", "Coût FCFA, technicien", "Durée, statut", "Rapport HTML"], ORANGE),
    ("Contrat", ["Fournisseur, montant", "Dates début/fin", "Renouvellement", "Statut actif/expiré"], PURPLE),
    ("Alerte", ["Type garantie/contrat", "Date échéance", "Message", "État actif/traité"], RED),
]

for i, (title, items, color) in enumerate(model_cards):
    x = Inches(0.4) + i * Inches(2.55)
    add_card(slide, title, items, x, Inches(1.8), Inches(2.35), Inches(4.5), color)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — DASHBOARD OWL
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, LIGHT_GRAY)
add_section_header(slide, "06", "Dashboard OWL", "Pilotage en temps réel")
add_footer(slide)

# KPI Cards
kpis = [
    ("Total Équipements", "Comptage global du parc", BLUE),
    ("Interventions en cours", "Planifiées + en cours", ORANGE),
    ("En maintenance", "Équipements en réparation", RED),
    ("Contrats expirants", "Échéance 90 jours", GREEN),
]

for i, (title, desc, color) in enumerate(kpis):
    x = Inches(0.6) + i * Inches(3.15)
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(2.85), Inches(1.8), WHITE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(2.85), Inches(0.07), color)

    # KPI icon placeholder (circle)
    icon = add_shape(slide, MSO_SHAPE.OVAL, x + Inches(1.0), Inches(2.0), Inches(0.85), Inches(0.85), color)
    add_text_box(slide, "KPI", x + Inches(1.0), Inches(2.1), Inches(0.85), Inches(0.65),
                 size=Pt(14), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, title, x + Inches(0.2), Inches(3.0), Inches(2.45), Inches(0.4),
                 size=Pt(16), color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x + Inches(0.2), Inches(3.4), Inches(2.45), Inches(0.4),
                 size=Pt(12), color=SUBTLE, alignment=PP_ALIGN.CENTER)

# Graphique section
add_text_box(slide, "Graphique : Répartition par catégorie", Inches(0.6), Inches(3.9), Inches(12), Inches(0.5),
             size=Pt(18), color=DARK_TEXT, bold=True)

chart_items = ["Ordinateurs (6)", "Serveurs (3)", "Réseau (4)", "Périphériques (2)"]
chart_colors = [BLUE, GREEN, ORANGE, PURPLE]
bar_max = Inches(2.5)
for i, (label, color) in enumerate(zip(chart_items, chart_colors)):
    x = Inches(0.8) + i * Inches(3.1)
    height = [Inches(1.8), Inches(1.2), Inches(1.5), Inches(0.8)][i]
    y = Inches(4.6) + (Inches(1.8) - height)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.5), y, Inches(1.6), height, color)
    add_text_box(slide, label, x, Inches(6.5), Inches(2.6), Inches(0.4),
                 size=Pt(13), color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — RAPPORTS & EXCEL
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, LIGHT_GRAY)
add_section_header(slide, "07", "Rapports PDF & Exports Excel", "Génération automatique de documents")
add_footer(slide)

# PDF column
add_card(slide, "Rapports PDF (QWeb)", [
    "Fiche équipement avec photo",
    "Rapport d'intervention",
    "Fiche contrat fournisseur",
    "Bouton dédié dans chaque vue",
], Inches(0.6), Inches(1.8), Inches(5.5), Inches(4.5), BLUE)

# Excel column
add_card(slide, "Exports Excel (xlsxwriter)", [
    "Inventaire complet du parc",
    "Coûts maintenance (équipement/mois)",
    "Contrats expirants sous 60 jours",
    "Couleurs conditionnelles :",
], Inches(6.6), Inches(1.8), Inches(6), Inches(4.5), GREEN)

# Color legend
legend_y = Inches(5.2)
colors_legend = [("Rouge  ≤ 15 jours restants", RED),
                 ("Orange  ≤ 30 jours restants", ORANGE),
                 ("Vert  > 30 jours restants", GREEN)]
for i, (text, color) in enumerate(colors_legend):
    x = Inches(6.8) + i * Inches(1.9)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, legend_y, Inches(0.25), Inches(0.25), color)
    add_text_box(slide, text, x + Inches(0.3), legend_y - Inches(0.05), Inches(1.6), Inches(0.3),
                 size=Pt(11), color=DARK_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — IMPORT & ALERTES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, LIGHT_GRAY)
add_section_header(slide, "08", "Import CSV & Alertes Automatiques", "Automatisation et productivité")
add_footer(slide)

add_card(slide, "Import CSV", [
    "Assistant Wizard upload fichier",
    "Validation avant import",
    "Création en masse équipements",
    "Gain de temps considérable",
], Inches(0.6), Inches(1.8), Inches(5.5), Inches(4.5), BLUE)

add_card(slide, "Alertes Automatiques", [
    "Scan quotidien (crons programmés)",
    "Garanties arrivant à expiration",
    "Contrats arrivant à échéance",
    "Création auto des alertes",
    "Notification dans le fil discussions",
], Inches(6.6), Inches(1.8), Inches(6), Inches(4.5), RED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — SÉCURITÉ
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, LIGHT_GRAY)
add_section_header(slide, "09", "Sécurité & Accès", "Deux profils distincts")
add_footer(slide)

# Technicien card
add_card(slide, "IT Technicien", [
    "Lecture équipements et contrats",
    "Création de ses interventions",
    "Dashboard et alertes",
    "Pas de suppression",
], Inches(0.8), Inches(1.8), Inches(5.3), Inches(4.5), GREEN)

# Manager card
add_card(slide, "IT Manager", [
    "Accès complet (CRUD)",
    "Exports Excel et rapports PDF",
    "Configuration du module",
    "Gestion des utilisateurs",
    "Toutes les fonctionnalités",
], Inches(6.8), Inches(1.8), Inches(5.3), Inches(4.5), BLUE)

# Security badges
badges = [("ACL", "ir.model.access.csv", BLUE), ("Rules", "Record Rules", GREEN), ("Groups", "2 Groupes", ORANGE)]
for i, (label, desc, color) in enumerate(badges):
    x = Inches(0.8) + i * Inches(2.0)
    badge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(6.5), Inches(1.8), Inches(0.4), color)
    add_text_box(slide, f"{label}  |  {desc}", x, Inches(6.5), Inches(1.8), Inches(0.4),
                 size=Pt(10), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — DONNÉES DÉMO
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, LIGHT_GRAY)
add_section_header(slide, "10", "Données de Démonstration", "Jeu de données réaliste")
add_footer(slide)

demo_cards = [
    ("15 Équipements", ["6 Ordinateurs", "3 Serveurs", "4 Réseau", "2 Périphériques"], BLUE),
    ("5 Contrats", ["Maintenance serveurs", "Licences Microsoft", "Cloud AWS", "Kaspersky", "Réseau"], GREEN),
    ("7 Interventions", ["4 Préventives", "2 Correctives", "1 Annulée"], ORANGE),
    ("4 Affectations", ["Historique complet", "Traçabilité totale"], PURPLE),
    ("3 Alertes", ["2 Garantie", "1 Contrat"], RED),
    ("7 Employés", ["DSI, Commercial, Support", "RH, DAF"], YELLOW),
]

for i, (title, items, color) in enumerate(demo_cards):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + col * Inches(4.1)
    y = Inches(1.8) + row * Inches(2.6)
    add_card(slide, title, items, x, y, Inches(3.8), Inches(2.3), color)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — DÉMO EN DIRECT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, LIGHT_GRAY)
add_section_header(slide, "11", "Démonstration en Direct", "8 cas d'usage concrets")
add_footer(slide)

demos = [
    ("01", "Création d'un équipement", "Nouveau PC + affectation employé", BLUE),
    ("02", "Planification intervention", "Préventive sur serveur", GREEN),
    ("03", "Tableau de bord", "KPI et graphique catégorie", ORANGE),
    ("04", "Rapport PDF", "Fiche équipement téléchargée", PURPLE),
    ("05", "Export Excel", "Inventaire avec couleurs", RED),
    ("06", "Import CSV", "Équipements en masse", BLUE_DARK),
    ("07", "Scan alertes", "Garanties et contrats", GREEN),
    ("08", "Droits utilisateurs", "Technicien vs Manager", YELLOW),
]

for i, (num, title, desc, color) in enumerate(demos):
    col = i % 4
    row = i // 4
    x = Inches(0.4) + col * Inches(3.2)
    y = Inches(1.8) + row * Inches(2.6)

    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.95), Inches(2.2), WHITE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(2.95), Inches(0.07), color)

    # Big number
    add_text_box(slide, num, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.6),
                 size=Pt(28), color=color, bold=True)
    # Title
    add_text_box(slide, title, x + Inches(0.2), y + Inches(0.8), Inches(2.55), Inches(0.4),
                 size=Pt(16), color=DARK_TEXT, bold=True)
    # Description
    add_text_box(slide, desc, x + Inches(0.2), y + Inches(1.3), Inches(2.55), Inches(0.5),
                 size=Pt(12), color=SUBTLE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, LIGHT_GRAY)
add_section_header(slide, "12", "Conclusion", "Bilan du projet")
add_footer(slide)

# Checkmark items
conclusion_items = [
    "Module Odoo 18 complet et fonctionnel répondant au cahier des charges TECHPARK CI",
    "9 fonctionnalités livrées : équipements, interventions, contrats, alertes, imports, exports, PDF, dashboard, sécurité",
    "Architecture modulaire et extensible — prêt pour la production",
    "Données démo réalistes pour présentation et tests",
    "Code versionné sur GitHub (15+ commits)",
]

for i, item in enumerate(conclusion_items):
    y = Inches(1.8) + Inches(0.75) * i
    # Checkmark circle
    circle = add_shape(slide, MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.4), Inches(0.4), GREEN)
    add_text_box(slide, "✓", Inches(0.6), y, Inches(0.4), Inches(0.4),
                 size=Pt(18), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, item, Inches(1.2), y + Inches(0.02), Inches(11.5), Inches(0.4),
                 size=Pt(17), color=DARK_TEXT)

# Next steps box
next_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0), WHITE)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.8), Inches(0.07), Inches(1.0), BLUE)
add_text_box(slide, "Prochaines étapes : Mise en production TECHPARK CI  |  Tickets support  |  Application mobile  |  QR code équipements",
             Inches(1.2), Inches(5.95), Inches(11), Inches(0.6),
             size=Pt(16), color=BLUE, bold=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — MERCI
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, PRIMARY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15), BLUE)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), prs.slide_width, Inches(0.15), BLUE)

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.0), Inches(5.3), Inches(0.06), BLUE)

add_text_box(slide, "Merci pour votre attention", Inches(1), Inches(1.8), Inches(11.3), Inches(1),
             size=Pt(48), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, "Silué Yacouba", Inches(1), Inches(3.3), Inches(11.3), Inches(0.6),
             size=Pt(28), color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, "Questions & Réponses", Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
             size=Pt(20), color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, "IT Parc TECHPARK CI  |  Odoo 18", Inches(1), Inches(5.2), Inches(11.3), Inches(0.4),
             size=Pt(14), color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── SAVE ──
path = os.path.expanduser("~/Desktop/IT_Parc_Presentation_v3.pptx")
prs.save(path)
print(f"PowerPoint enrichi créé : {path}")
print(f"Total slides : {len(prs.slides)}")
